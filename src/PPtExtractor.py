from pptx import Presentation


class PPtExtractor:
    def __init__(sefl):
        pass
    
    def extract(self, file_path:str) -> str:
        ppt = Presentation(file_path)
        text = []

        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return "\n".join(text).strip()
